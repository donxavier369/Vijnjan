import boto3
import os
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from django.conf import settings

def generate_s3_url(bucket_name, object_key):
    s3 = boto3.client('s3')
    return s3.generate_presigned_url(
        'get_object',
        Params={'Bucket': bucket_name, 'Key': object_key}
    )


def ppt_to_pdf(ppt_path, pdf_path):
    prs = Presentation(ppt_path)
    c = canvas.Canvas(pdf_path, pagesize=letter)
    width, height = letter

    for slide in prs.slides:
        # For simplicity, we just use the slide number as content here.
        # You might want to render slide content using python-pptx or other libraries.
        c.drawString(100, height - 100, f"Slide {prs.slides.index(slide) + 1}")
        c.showPage()

    c.save()
